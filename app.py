#pip install Flask python-docx Pillow pandas openpyxl PyPDF2 pdfplumber
#pip install virtualenv
#pip install textract

from flask import Flask, request, render_template, send_from_directory
import pdfplumber
import os
from docx import Document
import openpyxl
from pptx import Presentation

app = Flask(__name__)

UPLOAD_FOLDER = 'uploads'
TEXT_FOLDER = 'text_files'
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'xlsx', 'pptx'}

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
if not os.path.exists(TEXT_FOLDER):
    os.makedirs(TEXT_FOLDER)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        file = request.files['file']
        if file and allowed_file(file.filename):
            file_path = os.path.join(UPLOAD_FOLDER, file.filename)
            file.save(file_path)
            text = extract_text(file_path, file.filename)
            txt_filename = os.path.splitext(file.filename)[0] + '.txt'
            txt_path = os.path.join(TEXT_FOLDER, txt_filename)
            with open(txt_path, 'w', encoding='utf-8') as txt_file:
                txt_file.write(text)
            return render_template('download.html', txt_filename=txt_filename, text=text)
    return render_template('upload.html')

def extract_text(file_path, filename):
    extension = filename.rsplit('.', 1)[1].lower()
    text = ''
    if extension == 'pdf':
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ''
    elif extension == 'docx':
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + '\n'
    elif extension == 'xlsx':
        wb = openpyxl.load_workbook(file_path)
        sheet = wb.active
        for row in sheet.iter_rows(values_only=True):
            text += ' '.join(str(cell) for cell in row if cell is not None) + '\n'
    elif extension == 'pptx':
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + '\n'
    return text

@app.route('/download/<filename>')
def download(filename):
    return send_from_directory(TEXT_FOLDER, filename, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
